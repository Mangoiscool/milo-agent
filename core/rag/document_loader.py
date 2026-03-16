"""文档加载器

支持多种文档格式的加载和解析：
- PDF
- Markdown
- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)
- 纯文本
- 图像（OCR 或直接描述）
"""

import re
from abc import ABC, abstractmethod
from pathlib import Path

from .base import Document, DocumentType, detect_document_type


class BaseDocumentLoader(ABC):
    """文档加载器基类"""

    @abstractmethod
    def load(self, file_path: Path | str) -> list[Document]:
        """加载文档，返回文档列表（可能包含多页）"""
        pass

    @abstractmethod
    def supports(self, doc_type: DocumentType) -> bool:
        """检查是否支持该文档类型"""
        pass


class TextLoader(BaseDocumentLoader):
    """纯文本加载器"""

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()

        return [Document(
            content=content,
            source=str(path),
            doc_type=DocumentType.TEXT
        )]

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.TEXT


class MarkdownLoader(BaseDocumentLoader):
    """Markdown 文档加载器"""

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()

        return [Document(
            content=content,
            source=str(path),
            doc_type=DocumentType.MARKDOWN
        )]

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.MARKDOWN


class PDFLoader(BaseDocumentLoader):
    """PDF 文档加载器"""

    def __init__(self, extract_images: bool = False):
        """
        Args:
            extract_images: 是否提取 PDF 中的图片
        """
        self.extract_images = extract_images
        self._pypdf = None

    def _ensure_pypdf(self):
        """延迟导入 pypdf"""
        if self._pypdf is None:
            try:
                from pypdf import PdfReader
                self._pypdf = PdfReader
            except ImportError:
                raise ImportError(
                    "pypdf is required for PDF loading. "
                    "Install with: pip install pypdf"
                )
        return self._pypdf

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)
        PdfReader = self._ensure_pypdf()

        reader = PdfReader(str(path))
        documents = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                documents.append(Document(
                    content=text,
                    source=str(path),
                    doc_type=DocumentType.PDF,
                    page=i + 1
                ))

        return documents

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.PDF


class WordLoader(BaseDocumentLoader):
    """Word 文档加载器"""

    def __init__(self):
        self._docx = None

    def _ensure_docx(self):
        """延迟导入 python-docx"""
        if self._docx is None:
            try:
                from docx import Document as DocxDocument
                self._docx = DocxDocument
            except ImportError:
                raise ImportError(
                    "python-docx is required for Word loading. "
                    "Install with: pip install python-docx"
                )
        return self._docx

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)
        DocxDocument = self._ensure_docx()

        doc = DocxDocument(str(path))

        # 提取所有段落文本
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # 也提取表格内容
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                paragraphs.append("\n[表格]\n" + "\n".join(table_text))

        content = "\n\n".join(paragraphs)

        return [Document(
            content=content,
            source=str(path),
            doc_type=DocumentType.WORD
        )]

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.WORD


class ExcelLoader(BaseDocumentLoader):
    """Excel 文档加载器"""

    def __init__(self, sheet_name: str | None = None):
        """
        Args:
            sheet_name: 指定工作表名称，None 表示读取所有工作表
        """
        self.sheet_name = sheet_name
        self._openpyxl = None

    def _ensure_openpyxl(self):
        """延迟导入 openpyxl"""
        if self._openpyxl is None:
            try:
                import openpyxl
                self._openpyxl = openpyxl
            except ImportError:
                raise ImportError(
                    "openpyxl is required for Excel loading. "
                    "Install with: pip install openpyxl"
                )
        return self._openpyxl

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)
        openpyxl = self._ensure_openpyxl()

        workbook = openpyxl.load_workbook(str(path), data_only=True)
        documents = []

        sheets = [self.sheet_name] if self.sheet_name else workbook.sheetnames

        for sheet_name in sheets:
            if sheet_name not in workbook.sheetnames:
                continue

            sheet = workbook[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                # 过滤空行
                row_text = " | ".join(
                    str(cell) if cell is not None else ""
                    for cell in row
                )
                if row_text.strip(" |"):
                    rows.append(row_text)

            if rows:
                content = f"# 工作表: {sheet_name}\n\n" + "\n".join(rows)
                documents.append(Document(
                    content=content,
                    source=str(path),
                    doc_type=DocumentType.EXCEL,
                    metadata={"sheet_name": sheet_name}
                ))

        return documents

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.EXCEL


class PowerPointLoader(BaseDocumentLoader):
    """PowerPoint 文档加载器"""

    def __init__(self):
        self._pptx = None

    def _ensure_pptx(self):
        """延迟导入 python-pptx"""
        if self._pptx is None:
            try:
                from pptx import Presentation
                self._pptx = Presentation
            except ImportError:
                raise ImportError(
                    "python-pptx is required for PowerPoint loading. "
                    "Install with: pip install python-pptx"
                )
        return self._pptx

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)
        Presentation = self._ensure_pptx()

        prs = Presentation(str(path))
        documents = []

        for i, slide in enumerate(prs.slides):
            slide_content = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

            if slide_content:
                content = f"## 幻灯片 {i + 1}\n\n" + "\n\n".join(slide_content)
                documents.append(Document(
                    content=content,
                    source=str(path),
                    doc_type=DocumentType.POWERPOINT,
                    page=i + 1
                ))

        return documents

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.POWERPOINT


class ImageLoader(BaseDocumentLoader):
    """
    图像加载器

    支持两种模式：
    1. OCR 文字提取（需要 Tesseract 或其他 OCR 引擎）
    2. 直接返回图像描述占位符（等待多模态模型处理）
    """

    def __init__(self, use_ocr: bool = False):
        self.use_ocr = use_ocr

    def load(self, file_path: Path | str) -> list[Document]:
        path = Path(file_path)

        if self.use_ocr:
            content = self._ocr_image(path)
        else:
            # 返回占位符，等待后续多模态处理
            content = f"[图像文件: {path.name}]"

        return [Document(
            content=content,
            source=str(path),
            doc_type=DocumentType.IMAGE
        )]

    def _ocr_image(self, path: Path) -> str:
        """OCR 提取图像文字（需要额外依赖）"""
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(path)
            text = pytesseract.image_to_string(image, lang="chi_sim+eng")
            return text.strip()
        except ImportError:
            return f"[图像文件: {path.name} - OCR 不可用]"

    def supports(self, doc_type: DocumentType) -> bool:
        return doc_type == DocumentType.IMAGE


class DocumentLoaderRegistry:
    """文档加载器注册表"""

    def __init__(self):
        self._loaders: list[BaseDocumentLoader] = []

    def register(self, loader: BaseDocumentLoader):
        """注册加载器"""
        self._loaders.append(loader)

    def get_loader(self, doc_type: DocumentType) -> BaseDocumentLoader | None:
        """获取支持的加载器"""
        for loader in self._loaders:
            if loader.supports(doc_type):
                return loader
        return None

    def load(self, file_path: Path | str) -> list[Document]:
        """
        自动检测文件类型并加载

        Args:
            file_path: 文件路径

        Returns:
            文档列表
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        doc_type = detect_document_type(path)
        loader = self.get_loader(doc_type)

        if loader is None:
            raise ValueError(f"Unsupported document type: {doc_type.value}")

        return loader.load(path)

    def load_directory(
        self,
        directory: Path | str,
        recursive: bool = True,
        extensions: list[str] | None = None
    ) -> list[Document]:
        """
        加载目录下的所有文档

        Args:
            directory: 目录路径
            recursive: 是否递归子目录
            extensions: 文件扩展名过滤

        Returns:
            文档列表
        """
        dir_path = Path(directory)
        documents = []

        if recursive:
            pattern = "**/*"
        else:
            pattern = "*"

        for file_path in dir_path.glob(pattern):
            if not file_path.is_file():
                continue

            # 检查扩展名
            if extensions and file_path.suffix.lower() not in extensions:
                continue

            # 检查是否支持
            doc_type = detect_document_type(file_path)
            if doc_type == DocumentType.UNKNOWN:
                continue

            try:
                docs = self.load(file_path)
                documents.extend(docs)
            except Exception as e:
                print(f"Warning: Failed to load {file_path}: {e}")

        return documents


def create_default_registry() -> DocumentLoaderRegistry:
    """创建默认的文档加载器注册表"""
    registry = DocumentLoaderRegistry()
    registry.register(TextLoader())
    registry.register(MarkdownLoader())
    registry.register(PDFLoader())
    registry.register(WordLoader())
    registry.register(ExcelLoader())
    registry.register(PowerPointLoader())
    registry.register(ImageLoader())
    return registry